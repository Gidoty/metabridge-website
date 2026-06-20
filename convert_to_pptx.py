#!/usr/bin/env python3
"""Convert MetaBridge slide deck markdown files to formatted PowerPoint presentations.

Replicates the original MetaBridge visual style:
  - Dark navy background (#0D2137)
  - Cyan accent bar and labels (#00C2D4)
  - Left-side module info panel
  - Card-based content with coloured left borders
  - Course-specific accent colours
"""

import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette (matched from original MetaBridge PPTX) ──────────────────
BG_MAIN   = RGBColor(0x0D, 0x21, 0x37)  # main dark navy
BG_DARK   = RGBColor(0x0A, 0x19, 0x29)  # darker navy (footer / cards)
BG_CARD   = RGBColor(0x12, 0x22, 0x36)  # card background
CYAN      = RGBColor(0x00, 0xC2, 0xD4)  # primary cyan
RED       = RGBColor(0xE7, 0x4C, 0x3C)  # red (cybersecurity)
GOLD      = RGBColor(0xF5, 0xA6, 0x23)  # gold (blockchain / crypto)
GREEN     = RGBColor(0x00, 0xB8, 0x94)  # green
PURPLE    = RGBColor(0x9B, 0x59, 0xB6)  # purple (AI / prompt)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHT  = RGBColor(0xF0, 0xF4, 0xF8)  # near-white body text
MUTED     = RGBColor(0x8B, 0xA5, 0xBE)  # muted blue-grey
WARNING   = RGBColor(0xFF, 0x80, 0x80)  # warning text (light red)

COURSE_ACCENT = {
    "cybersecurity":             RED,
    "data-analytics":            CYAN,
    "ai-prompt-engineering":     PURPLE,
    "blockchain-cryptocurrency": GOLD,
}

COURSE_LABELS = {
    "cybersecurity":             "Cybersecurity",
    "data-analytics":            "Data Analytics",
    "ai-prompt-engineering":     "AI & Prompt Engineering",
    "blockchain-cryptocurrency": "Blockchain & Cryptocurrency",
}

# ── Slide dimensions (standard widescreen 16:9) ─────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Layout constants ─────────────────────────────────────────────────────────
TOP_BAR_H   = Inches(0.07)          # thin accent bar at top
SIDEBAR_W   = Inches(4.7)           # left module info area
DIVIDER_W   = Inches(0.06)          # thin cyan divider
CONTENT_X   = SIDEBAR_W + DIVIDER_W # right content area start
CONTENT_W   = SLIDE_W - CONTENT_X   # right content area width
FOOTER_H    = Inches(0.32)          # bottom footer bar
FOOTER_Y    = SLIDE_H - FOOTER_H


# ── Core drawing helpers ─────────────────────────────────────────────────────

def set_bg(slide, colour):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_rect(slide, x, y, w, h, fill_colour, line_colour=None, line_width=Pt(0)):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    if line_colour:
        shape.line.color.rgb = line_colour
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, colour=WHITE,
                align=PP_ALIGN.LEFT, wrap=True, italic=False, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    run.font.name = font_name
    return txBox


def add_textbox_multiline(slide, lines, x, y, w, h,
                           font_size=11, colour=NEAR_WHT, font_name='Calibri'):
    """Add a textbox containing multiple lines with consistent formatting."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line_text, line_colour, line_size, line_bold in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line_text
        run.font.size = Pt(line_size)
        run.font.bold = line_bold
        run.font.color.rgb = line_colour
        run.font.name = font_name
    return txBox


# ── Card helper ──────────────────────────────────────────────────────────────

def add_card(slide, x, y, w, h, border_colour, title, desc, icon=''):
    """Card with coloured left border, title, description."""
    # Card background
    add_rect(slide, x, y, w, h, BG_CARD)
    # Left colour strip
    add_rect(slide, x, y, Inches(0.055), h, border_colour)

    inner_x = x + Inches(0.12)
    inner_w = w - Inches(0.18)

    if icon:
        add_textbox(slide, icon, inner_x, y + Inches(0.08), Inches(0.4), Inches(0.4),
                    font_size=15, colour=WHITE)
        title_x = inner_x + Inches(0.42)
        title_w = inner_w - Inches(0.42)
    else:
        title_x = inner_x
        title_w = inner_w

    add_textbox(slide, title, title_x, y + Inches(0.1), title_w, Inches(0.4),
                font_size=11, bold=True, colour=border_colour, wrap=True)

    desc_y = y + Inches(0.5) if title else y + Inches(0.1)
    desc_h = h - Inches(0.58)
    if desc:
        add_textbox(slide, desc, inner_x, desc_y, inner_w, max(desc_h, Inches(0.3)),
                    font_size=9.5, colour=NEAR_WHT, wrap=True)


# ── Sidebar helpers ──────────────────────────────────────────────────────────

def draw_sidebar(slide, mod_label, mod_title, course_name, accent):
    """Draw the left module info panel on a content slide."""
    # Top accent bar (full width)
    add_rect(slide, 0, 0, SLIDE_W, TOP_BAR_H, accent)

    # Cyan vertical divider between sidebar and content
    add_rect(slide, SIDEBAR_W, TOP_BAR_H, DIVIDER_W,
             SLIDE_H - TOP_BAR_H - FOOTER_H, accent)

    # Footer bar (full width)
    add_rect(slide, 0, FOOTER_Y, SLIDE_W, FOOTER_H, BG_DARK)

    # Footer text
    add_textbox(slide, f"MetaBridge Academy  |  {course_name}  |  2025",
                Inches(0.2), FOOTER_Y, Inches(9), FOOTER_H,
                font_size=8.5, colour=MUTED, align=PP_ALIGN.LEFT)
    add_textbox(slide, "metabridgeacademy.com",
                SLIDE_W - Inches(2.2), FOOTER_Y, Inches(2.0), FOOTER_H,
                font_size=8.5, colour=accent, align=PP_ALIGN.RIGHT)

    # Module badge
    if mod_label:
        add_textbox(slide, mod_label,
                    Inches(0.25), TOP_BAR_H + Inches(0.15), SIDEBAR_W - Inches(0.3), Inches(0.45),
                    font_size=15, bold=True, colour=accent)

    # Slide/module title (large, wrapping)
    if mod_title:
        add_textbox(slide, mod_title,
                    Inches(0.25), TOP_BAR_H + Inches(0.65), SIDEBAR_W - Inches(0.3), Inches(4.5),
                    font_size=22, bold=True, colour=WHITE, wrap=True)

    # MetaBridge Academy watermark at bottom of sidebar
    add_textbox(slide, "MetaBridge Academy",
                Inches(0.25), FOOTER_Y - Inches(0.4), SIDEBAR_W - Inches(0.3), Inches(0.35),
                font_size=9, colour=MUTED, italic=True)


# ── Slide type builders ──────────────────────────────────────────────────────

def make_title_slide(prs, course_name, accent):
    """Cover slide — two-panel with stats boxes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_MAIN)

    # Top accent bar
    add_rect(slide, 0, 0, SLIDE_W, TOP_BAR_H, accent)

    # Left panel (darker navy)
    add_rect(slide, 0, TOP_BAR_H, SIDEBAR_W, SLIDE_H - TOP_BAR_H, BG_DARK)

    # Cyan divider
    add_rect(slide, SIDEBAR_W, TOP_BAR_H, DIVIDER_W, SLIDE_H - TOP_BAR_H, accent)

    # Left panel content
    add_textbox(slide, "MetaBridge", Inches(0.3), Inches(0.4), Inches(4.2), Inches(0.55),
                font_size=28, bold=True, colour=accent)
    add_textbox(slide, "Academy", Inches(0.3), Inches(0.9), Inches(4.2), Inches(0.5),
                font_size=28, bold=False, colour=WHITE)

    add_textbox(slide, course_name,
                Inches(0.3), Inches(2.0), Inches(4.2), Inches(1.8),
                font_size=26, bold=True, colour=WHITE, wrap=True)

    add_textbox(slide, "Comprehensive Course\nSlide Deck",
                Inches(0.3), Inches(3.9), Inches(4.2), Inches(0.9),
                font_size=14, colour=MUTED, wrap=True)

    # Footer bottom panel
    add_rect(slide, 0, FOOTER_Y, SLIDE_W, FOOTER_H, BG_DARK)
    add_textbox(slide, "Cybersecurity  |  Data Analytics  |  AI & Prompt Engineering  |  Blockchain & Cryptocurrency",
                SIDEBAR_W + DIVIDER_W + Inches(0.2), FOOTER_Y, CONTENT_W - Inches(0.3), FOOTER_H,
                font_size=8, colour=MUTED, align=PP_ALIGN.CENTER)

    # Right panel stat boxes
    stat_y = Inches(1.2)
    stats = [
        (course_name.split()[0], "Course", accent),
        ("12 Modules", "of content", GOLD if accent != GOLD else CYAN),
        ("72+ Hours", "of training", GREEN),
        ("Blockchain", "Certificate", MUTED),
    ]
    sx = CONTENT_X + Inches(0.4)
    sw = (CONTENT_W - Inches(0.8)) / 2 - Inches(0.15)
    sy = [Inches(1.2), Inches(1.2), Inches(3.0), Inches(3.0)]
    sxs = [sx, sx + sw + Inches(0.3), sx, sx + sw + Inches(0.3)]
    stat_colors = [accent, GOLD if accent != GOLD else GREEN, GREEN, MUTED]

    for i, (big, small, sc) in enumerate(stats):
        add_rect(slide, sxs[i], sy[i], sw, Inches(1.55), BG_CARD)
        add_rect(slide, sxs[i], sy[i], Inches(0.055), Inches(1.55), sc)
        add_textbox(slide, big, sxs[i] + Inches(0.15), sy[i] + Inches(0.15), sw - Inches(0.2), Inches(0.7),
                    font_size=20, bold=True, colour=sc, wrap=True)
        add_textbox(slide, small, sxs[i] + Inches(0.15), sy[i] + Inches(0.85), sw - Inches(0.2), Inches(0.5),
                    font_size=11, colour=NEAR_WHT)

    # Stats strip at bottom of right panel
    strip_y = FOOTER_Y - Inches(0.55)
    add_rect(slide, CONTENT_X, strip_y, CONTENT_W, Inches(0.5), BG_CARD)
    add_textbox(slide, "12 Modules  ·  72+ Hours  ·  Hands-On Labs  ·  CompTIA S+ Aligned  ·  Blockchain Certificate",
                CONTENT_X + Inches(0.2), strip_y, CONTENT_W - Inches(0.4), Inches(0.5),
                font_size=9.5, colour=MUTED, align=PP_ALIGN.CENTER)
    return slide


def make_module_opener(prs, mod_label, title, content_lines, notes, course_name, accent):
    """Module/section opener slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_MAIN)
    draw_sidebar(slide, '', '', course_name, accent)

    # Large ghost module number (decorative)
    num_match = re.search(r'\d+', mod_label) if mod_label else None
    if num_match:
        ghost_num = num_match.group()
        add_textbox(slide, ghost_num,
                    Inches(0.0), Inches(0.0), SIDEBAR_W, Inches(5.5),
                    font_size=160, bold=True,
                    colour=RGBColor(0x1A, 0x33, 0x50), align=PP_ALIGN.LEFT)

    # Module label
    add_textbox(slide, mod_label,
                Inches(0.3), Inches(1.2), SIDEBAR_W - Inches(0.4), Inches(0.6),
                font_size=16, bold=True, colour=accent)

    # Title (large, prominent)
    add_textbox(slide, title,
                Inches(0.3), Inches(2.0), SIDEBAR_W - Inches(0.4), Inches(4.0),
                font_size=30, bold=True, colour=WHITE, wrap=True)

    # Right side: content as bullet list
    content_x = CONTENT_X + Inches(0.35)
    content_w = CONTENT_W - Inches(0.5)
    y = TOP_BAR_H + Inches(0.5)

    txBox = slide.shapes.add_textbox(content_x, y, content_w, Inches(5.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True

    for line in content_lines:
        stripped = line.strip()
        if not stripped:
            continue
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        run = p.add_run()
        # Header lines (ALL CAPS)
        if re.match(r'^[A-Z][A-Z0-9\s&·\-–—:()/]+$', stripped) and len(stripped) > 4:
            run.text = stripped
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = accent
        elif stripped.startswith('- ') or stripped.startswith('→ ') or stripped.startswith('• '):
            run.text = '  •  ' + stripped[2:]
            run.font.size = Pt(11)
            run.font.color.rgb = NEAR_WHT
        elif stripped.startswith('LEVEL:'):
            run.text = stripped
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = GOLD
        else:
            run.text = stripped
            run.font.size = Pt(11)
            run.font.color.rgb = NEAR_WHT
        run.font.name = 'Calibri'

    if notes.strip():
        slide.notes_slide.notes_text_frame.text = notes.strip()
    return slide


def make_content_slide(prs, mod_label, mod_title, slide_title, context,
                       content_lines, notes, course_name, accent):
    """Standard content slide with intelligent layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_MAIN)
    draw_sidebar(slide, mod_label, mod_title, course_name, accent)

    cx = CONTENT_X + Inches(0.3)
    cw = CONTENT_W - Inches(0.45)

    # Context line
    if context:
        add_textbox(slide, context,
                    cx, TOP_BAR_H + Inches(0.12), cw, Inches(0.38),
                    font_size=9.5, colour=MUTED, italic=True, wrap=True)

    # Slide title in content area
    add_textbox(slide, slide_title,
                cx, TOP_BAR_H + Inches(0.5), cw, Inches(0.65),
                font_size=18, bold=True, colour=WHITE)

    # Thin separator
    add_rect(slide, cx, TOP_BAR_H + Inches(1.18), cw, Inches(0.025), accent)

    # Content area
    content_top = TOP_BAR_H + Inches(1.25)
    content_h = FOOTER_Y - content_top - Inches(0.1)

    # Try card layout first
    cards = _extract_cards(content_lines)
    if cards and len(cards) >= 2:
        _draw_cards(slide, cx, content_top, cw, content_h, cards, accent)
    else:
        _draw_text_content(slide, cx, content_top, cw, content_h, content_lines, accent)

    if notes.strip():
        slide.notes_slide.notes_text_frame.text = notes.strip()
    return slide


def _extract_cards(content_lines):
    """Try to extract card items (title + description) from content lines."""
    cards = []
    i = 0
    lines = [l.strip() for l in content_lines]

    while i < len(lines):
        line = lines[i]
        if not line:
            i += 1
            continue

        # Pattern: ALL CAPS title (3–60 chars) followed by description lines
        if (re.match(r'^[A-Z][A-Z0-9\s&·\-–—:()/]+$', line)
                and 4 < len(line) < 65
                and not line.startswith('---')):
            desc_parts = []
            j = i + 1
            while j < len(lines) and j < i + 6:
                next_l = lines[j]
                if not next_l:
                    break
                # Stop if another ALL CAPS header is encountered
                if re.match(r'^[A-Z][A-Z0-9\s&·\-–—:()/]+$', next_l) and len(next_l) > 4:
                    break
                desc_parts.append(re.sub(r'^[-→•·]\s*', '', next_l))
                j += 1
            if desc_parts:
                cards.append({'title': line, 'desc': ' '.join(desc_parts[:3]), 'icon': ''})
                i = j
                continue

        # Pattern: → Item or • Item (icon-led cards)
        if line.startswith('→ ') or (line.startswith('• ') and len(line) > 4):
            prefix_len = 2
            title = line[prefix_len:].split('—')[0].split(':')[0].strip()
            rest = line[prefix_len:]
            desc_parts = [rest]
            j = i + 1
            while j < len(lines) and j < i + 3:
                next_l = lines[j]
                if not next_l or next_l.startswith('→') or next_l.startswith('•'):
                    break
                if next_l.startswith(' ') or next_l.startswith('\t'):
                    desc_parts.append(next_l.strip())
                j += 1
            cards.append({'title': title, 'desc': rest, 'icon': '→'})
            i = j
            continue

        i += 1

    return cards if len(cards) >= 2 else []


def _draw_cards(slide, cx, cy, cw, ch, cards, accent):
    """Draw cards in a grid layout."""
    n = len(cards)
    cols = 2 if n > 3 else min(n, 3)
    rows = (n + cols - 1) // cols

    gap_x = Inches(0.18)
    gap_y = Inches(0.16)
    card_w = (cw - gap_x * (cols - 1)) / cols
    card_h = min((ch - gap_y * (rows - 1)) / rows, Inches(1.9))

    # Cycle accent colours for visual interest
    border_colours = [accent, GOLD if accent != GOLD else GREEN,
                      GREEN if accent != GREEN else CYAN,
                      CYAN if accent != CYAN else GOLD,
                      RED if accent != RED else accent, accent]

    for idx, card in enumerate(cards[:rows * cols]):
        row = idx // cols
        col = idx % cols
        x = cx + col * (card_w + gap_x)
        y = cy + row * (card_h + gap_y)
        bc = border_colours[idx % len(border_colours)]
        add_card(slide, x, y, card_w, card_h, bc,
                 card['title'], card['desc'], card.get('icon', ''))


def _draw_text_content(slide, cx, cy, cw, ch, content_lines, accent):
    """Draw content as a formatted text block."""
    txBox = slide.shapes.add_textbox(cx, cy, cw, ch)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True

    for line in content_lines:
        stripped = line.strip()
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        run = p.add_run()
        run.font.name = 'Calibri'

        if not stripped:
            run.text = ''
            run.font.size = Pt(6)
            continue

        # ALL CAPS section header
        if re.match(r'^[A-Z][A-Z0-9\s&·\-–—:()/]+$', stripped) and len(stripped) > 4:
            run.text = stripped
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = accent
        # Bullet / arrow items
        elif (stripped.startswith('- ') or stripped.startswith('→ ')
              or stripped.startswith('• ') or stripped.startswith('· ')):
            run.text = '  •  ' + stripped[2:]
            run.font.size = Pt(10.5)
            run.font.color.rgb = NEAR_WHT
        # Numbered items
        elif re.match(r'^\d+[.)]\s', stripped):
            run.text = '  ' + stripped
            run.font.size = Pt(10.5)
            run.font.color.rgb = NEAR_WHT
        # Warning lines (⚠)
        elif '⚠' in stripped:
            run.text = stripped
            run.font.size = Pt(10)
            run.font.color.rgb = WARNING
        # Code / terminal lines
        elif (stripped.startswith('=') or stripped.startswith('$')
              or '()' in stripped or '->' in stripped
              or stripped.startswith('`')):
            run.text = stripped.replace('`', '')
            run.font.size = Pt(9.5)
            run.font.color.rgb = RGBColor(0x4A, 0xDE, 0x80)
            run.font.name = 'Courier New'
        # Sub-bullets (indented)
        elif line.startswith('  ') and stripped:
            run.text = '      ◦  ' + stripped
            run.font.size = Pt(10)
            run.font.color.rgb = MUTED
        else:
            # Handle inline **bold**
            parts = re.split(r'(\*\*[^*]+\*\*)', stripped)
            if len(parts) == 1:
                run.text = stripped
                run.font.size = Pt(10.5)
                run.font.color.rgb = NEAR_WHT
            else:
                # Multi-run paragraph
                run.text = ''
                for part in parts:
                    r = p.add_run()
                    r.font.name = 'Calibri'
                    r.font.size = Pt(10.5)
                    if part.startswith('**') and part.endswith('**'):
                        r.text = part[2:-2]
                        r.font.bold = True
                        r.font.color.rgb = WHITE
                    else:
                        r.text = part
                        r.font.color.rgb = NEAR_WHT


# ── Markdown parser ──────────────────────────────────────────────────────────

def parse_slide_deck(md_path):
    """Parse the slide deck markdown into a list of slide dicts."""
    with open(md_path, 'r', encoding='utf-8') as f:
        raw = f.read()

    # Split on ## SLIDE N markers
    blocks = re.split(r'\n(?=## SLIDE\s)', raw, flags=re.IGNORECASE)

    slides = []
    for block in blocks:
        block = block.strip()
        if not block:
            continue

        m = re.match(r'^## SLIDE\s+\d+[A-Za-z]?\s*[—–\-]+\s*(.+)', block, re.IGNORECASE)
        if not m:
            continue

        title = m.group(1).strip()

        # Extract content block (inside ``` ```)
        content_match = re.search(r'\*\*Content:\*\*\s*```[^\n]*\n([\s\S]+?)```', block)
        if content_match:
            content_lines = content_match.group(1).rstrip().split('\n')
        else:
            cm2 = re.search(r'\*\*Content:\*\*\s*\n([\s\S]+?)(?=\*\*Speaker|\*\*Visual|\Z)', block)
            content_lines = cm2.group(1).strip().split('\n') if cm2 else []

        # Extract visual note (used as context line)
        visual_m = re.search(r'\*\*Visual:\*\*\s*(.+?)(?=\n\n|\*\*Content|\*\*Speaker|\Z)', block, re.DOTALL)
        visual = visual_m.group(1).strip().replace('\n', ' ') if visual_m else ''

        # Extract speaker notes
        notes_m = re.search(r'\*\*Speaker notes?:\*\*\s*([\s\S]+?)(?=---|\Z)', block)
        notes = notes_m.group(1).strip() if notes_m else ''
        notes = re.sub(r'\*\*([^*]+)\*\*', r'\1', notes)

        slides.append({
            'title': title,
            'content': content_lines,
            'visual': visual,
            'notes': notes,
        })

    return slides


def _get_module_label(slide_data, state):
    """Determine the MOD N label for a slide, tracking state across slides."""
    content_text = '\n'.join(slide_data['content'])
    title = slide_data['title']

    # Check title for module reference
    mod_in_title = re.search(r'MODULE\s+(\d+)', title, re.IGNORECASE)
    if mod_in_title:
        state['mod_num'] = int(mod_in_title.group(1))
        return f"MOD {state['mod_num']}"

    # Check content for module reference
    mod_in_content = re.search(r'MODULE\s+0?(\d+)', content_text, re.IGNORECASE)
    if mod_in_content:
        state['mod_num'] = int(mod_in_content.group(1))
        return f"MOD {state['mod_num']}"

    if state['mod_num']:
        return f"MOD {state['mod_num']}"
    return "MOD —"


def _get_module_title(slide_data, state):
    """Return the current module title, updating state on new module openers."""
    title = slide_data['title']
    content_text = '\n'.join(slide_data['content'])

    if 'OPENER' in title.upper() or re.search(r'MODULE\s+\d+\s*OPENER', title, re.IGNORECASE):
        state['mod_title'] = title.replace('MODULE ', 'Module ').split(':')[-1].strip()
    elif re.search(r'^MODULE\s+\d+\s*[:\-]', title, re.IGNORECASE):
        state['mod_title'] = re.sub(r'^MODULE\s+\d+\s*[:\-]\s*', '', title, flags=re.IGNORECASE).strip()

    return state.get('mod_title', '')


# ── Main converter ───────────────────────────────────────────────────────────

def convert_slide_deck(md_path, pptx_path, course_key):
    accent = COURSE_ACCENT.get(course_key, CYAN)
    course_name = COURSE_LABELS.get(course_key, course_key)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    make_title_slide(prs, course_name, accent)

    slides = parse_slide_deck(md_path)
    state = {'mod_num': 0, 'mod_title': ''}

    for sd in slides:
        title = sd['title']
        content = sd['content']
        notes = sd['notes']
        visual = sd['visual']

        # Skip generic TITLE SLIDE entry — we already create a custom cover slide
        if re.match(r'^TITLE\s+SLIDE$', title.strip(), re.IGNORECASE):
            continue

        mod_label = _get_module_label(sd, state)
        mod_title = _get_module_title(sd, state)

        # Detect module opener slides
        is_opener = ('OPENER' in title.upper()
                     or re.search(r'^MODULE\s+\d+\s*OPENER', title, re.IGNORECASE)
                     or (re.match(r'^MODULE\s+\d+\b', title, re.IGNORECASE) and len(content) < 15))

        if is_opener:
            make_module_opener(prs, mod_label, title, content, notes, course_name, accent)
        else:
            # Build a short context line from visual description (first sentence only)
            context = visual.split('.')[0][:120] if visual else ''
            make_content_slide(prs, mod_label, mod_title, title,
                               context, content, notes, course_name, accent)

    prs.save(pptx_path)
    print(f"  Saved: {pptx_path}")


# ── Entry point ──────────────────────────────────────────────────────────────

if __name__ == '__main__':
    base = '/home/user/metabridge-website/curriculum'
    out_dir = '/home/user/metabridge-website/exports'
    os.makedirs(out_dir, exist_ok=True)

    courses = [
        'cybersecurity',
        'data-analytics',
        'ai-prompt-engineering',
        'blockchain-cryptocurrency',
    ]

    labels = {
        "cybersecurity":             "Cybersecurity",
        "data-analytics":            "Data-Analytics",
        "ai-prompt-engineering":     "AI-and-Prompt-Engineering",
        "blockchain-cryptocurrency": "Blockchain-and-Cryptocurrency",
    }

    for course in courses:
        src = os.path.join(base, course, '03-slide-deck.md')
        if os.path.exists(src):
            course_label = labels[course]
            dst = os.path.join(out_dir, f"{course_label}-Slide-Deck.pptx")
            print(f"Converting: {course}/03-slide-deck.md")
            convert_slide_deck(src, dst, course)
        else:
            print(f"  MISSING: {src}")

    print("\nPowerPoint conversion complete.")
